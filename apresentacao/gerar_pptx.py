"""Monta a apresentacao em .pptx nativo: texto editavel, nao imagem.

So o histograma entra como figura, porque e um grafico. Todo o resto sao
caixas de texto e formas que abrem editaveis no PowerPoint.

Os numeros vem de numeros.json, gerado junto com a figura, para nao haver
valor digitado a mao que divirja do corpus.

Uso:
    .venv/bin/python apresentacao/gerar_pptx.py
"""

import json
import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

AQUI = pathlib.Path(__file__).parent
RAIZ = AQUI.parent
SAIDA = AQUI / "apresentacao.pptx"

# Arial em vez de Helvetica Neue: existe no Mac e no Windows, entao o arquivo
# nao troca de fonte na maquina de outra pessoa.
FONTE = "Arial"

PAPEL = RGBColor(0xFC, 0xFC, 0xFB)
TINTA = RGBColor(0x0B, 0x0B, 0x0B)
TINTA2 = RGBColor(0x52, 0x51, 0x4E)
MUDO = RGBColor(0x89, 0x87, 0x81)
LINHA = RGBColor(0xE1, 0xE0, 0xD9)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
NEG = RGBColor(0xE3, 0x49, 0x48)
NEU = RGBColor(0x89, 0x87, 0x81)
POS = RGBColor(0x39, 0x87, 0xE5)
NEG_FUNDO = RGBColor(0xFD, 0xEE, 0xEE)
POS_FUNDO = RGBColor(0xEE, 0xF4, 0xFD)

MARGEM = Inches(0.7)
LARGURA_UTIL = Inches(13.333) - 2 * MARGEM


# --------------------------------------------------------------- utilidades
def texto(slide, x, y, w, h, conteudo, tamanho: float = 12, cor=TINTA2, negrito=False,
          alinhamento=PP_ALIGN.LEFT, espacamento=1.25, ancora=MSO_ANCHOR.TOP):
    """Caixa de texto. `conteudo` aceita **negrito** e _italico_ por trecho."""
    caixa = slide.shapes.add_textbox(x, y, w, h)
    quadro = caixa.text_frame
    quadro.word_wrap = True
    quadro.vertical_anchor = ancora
    quadro.margin_left = quadro.margin_right = 0
    quadro.margin_top = quadro.margin_bottom = 0

    for indice, linha in enumerate(conteudo.split("\n")):
        paragrafo = quadro.paragraphs[0] if indice == 0 else quadro.add_paragraph()
        paragrafo.alignment = alinhamento
        paragrafo.line_spacing = espacamento
        if indice:
            paragrafo.space_before = Pt(6)
        _preencher(paragrafo, linha, tamanho, cor, negrito)
    return caixa


def _preencher(paragrafo, linha, tamanho: float, cor, negrito):
    """Quebra a linha em trechos normais, **negrito** e _italico_."""
    import re
    for pedaco in re.split(r"(\*\*[^*]+\*\*|_[^_]+_)", linha):
        if not pedaco:
            continue
        corrida = paragrafo.add_run()
        if pedaco.startswith("**"):
            corrida.text = pedaco[2:-2]
            corrida.font.bold = True
            corrida.font.color.rgb = TINTA if cor is TINTA2 else cor
        elif pedaco.startswith("_"):
            corrida.text = pedaco[1:-1]
            corrida.font.italic = True
            corrida.font.color.rgb = cor
        else:
            corrida.text = pedaco
            corrida.font.bold = negrito
            corrida.font.color.rgb = cor
        corrida.font.size = Pt(tamanho)
        corrida.font.name = FONTE


def forma(slide, tipo, x, y, w, h, fundo=None, borda=None, espessura=0.75):
    peca = slide.shapes.add_shape(tipo, x, y, w, h)
    if fundo is None:
        peca.fill.background()
    else:
        peca.fill.solid()
        peca.fill.fore_color.rgb = fundo
    if borda is None:
        peca.line.fill.background()
    else:
        peca.line.color.rgb = borda
        peca.line.width = Pt(espessura)
    peca.shadow.inherit = False
    peca.text_frame.text = ""
    return peca


def fundo_papel(slide):
    forma(slide, MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5), fundo=PAPEL)


def cabecalho(slide, titulo, subtitulo=None, fraco=None):
    y = Inches(0.55)
    caixa = slide.shapes.add_textbox(MARGEM, y, LARGURA_UTIL, Inches(0.6))
    quadro = caixa.text_frame
    quadro.word_wrap = True
    quadro.margin_left = quadro.margin_right = quadro.margin_top = quadro.margin_bottom = 0
    p = quadro.paragraphs[0]
    p.line_spacing = 1.1
    forte = p.add_run()
    forte.text = titulo
    forte.font.size, forte.font.bold, forte.font.name = Pt(27), True, FONTE
    forte.font.color.rgb = TINTA
    if fraco:
        leve = p.add_run()
        leve.text = fraco
        leve.font.size, leve.font.name = Pt(27), FONTE
        leve.font.color.rgb = MUDO
    if subtitulo:
        texto(slide, MARGEM, Inches(1.18), LARGURA_UTIL, Inches(0.5),
              subtitulo, tamanho=12.5, cor=MUDO, espacamento=1.3)


def rodape(slide, esquerda, numero):
    y = Inches(6.95)
    linha = slide.shapes.add_connector(1, MARGEM, y, MARGEM + LARGURA_UTIL, y)
    linha.line.color.rgb = LINHA
    linha.line.width = Pt(0.75)
    texto(slide, MARGEM, y + Inches(0.09), LARGURA_UTIL - Inches(0.5), Inches(0.3),
          esquerda, tamanho=8.5, cor=MUDO)
    texto(slide, MARGEM + LARGURA_UTIL - Inches(0.4), y + Inches(0.09),
          Inches(0.4), Inches(0.3), str(numero), tamanho=8.5, cor=MUDO,
          alinhamento=PP_ALIGN.RIGHT)


def cartao(slide, x, y, w, h, acento=False):
    """Caixa branca com borda. Com acento, ganha barra escura na esquerda."""
    forma(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fundo=BRANCO, borda=LINHA)
    if acento:
        forma(slide, MSO_SHAPE.RECTANGLE, x, y, Inches(0.05), h, fundo=TINTA)


def passo(slide, x, y, w, numero, titulo, corpo):
    diametro = Inches(0.34)
    bolha = forma(slide, MSO_SHAPE.OVAL, x, y, diametro, diametro, fundo=TINTA)
    quadro = bolha.text_frame
    quadro.margin_left = quadro.margin_right = 0
    quadro.margin_top = quadro.margin_bottom = 0
    quadro.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = quadro.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(numero)
    r.font.size, r.font.bold, r.font.name = Pt(12), True, FONTE
    r.font.color.rgb = PAPEL

    recuo = x + Inches(0.5)
    largura = w - Inches(0.5)
    texto(slide, recuo, y - Inches(0.02), largura, Inches(0.3), titulo,
          tamanho=14.5, cor=TINTA, negrito=True)
    texto(slide, recuo, y + Inches(0.28), largura, Inches(0.9), corpo,
          tamanho=11.5, cor=TINTA2, espacamento=1.3)


def barra_empilhada(slide, x, y, w, altura, fatias):
    """fatias = [(percentual, cor, rotulo ou None)]"""
    esquerda = x
    for pct, cor, rotulo in fatias:
        if pct <= 0:
            continue
        largura = Emu(int(w * pct / 100))
        peca = forma(slide, MSO_SHAPE.RECTANGLE, esquerda, y, largura, altura, fundo=cor)
        if rotulo:
            quadro = peca.text_frame
            quadro.margin_left = quadro.margin_right = 0
            quadro.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = quadro.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = rotulo
            r.font.size, r.font.bold, r.font.name = Pt(10), True, FONTE
            r.font.color.rgb = BRANCO
        esquerda += largura


# ------------------------------------------------------------------ slides
def capa(apresentacao, d):
    slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])
    fundo_papel(slide)

    forma(slide, MSO_SHAPE.RECTANGLE, MARGEM, Inches(2.15), Inches(1.5), Inches(0.07),
          fundo=POS)

    texto(slide, MARGEM, Inches(2.5), Inches(10.5), Inches(1.5),
          "Do que as pessoas falam\nquando falam da Argentina",
          tamanho=40, cor=TINTA, negrito=True, espacamento=1.12)

    texto(slide, MARGEM, Inches(4.35), Inches(9.5), Inches(0.9),
          "Análise de sentimento de posts de rede social com o dicionário ANEW\n"
          "de normas afetivas (Bradley & Lang, 1999)",
          tamanho=15, cor=MUDO, espacamento=1.4)

    y = Inches(5.6)
    linha = slide.shapes.add_connector(1, MARGEM, y, MARGEM + Inches(5.2), y)
    linha.line.color.rgb = LINHA
    linha.line.width = Pt(0.75)

    texto(slide, MARGEM, y + Inches(0.25), Inches(5.5), Inches(0.8),
          "Pedro Fardin\n(preencher o nome do colega)",
          tamanho=12.5, cor=TINTA2, espacamento=1.4)

    texto(slide, MARGEM + Inches(6.5), y + Inches(0.25), Inches(5.4), Inches(0.8),
          "Processamento de Linguagem Natural\n"
          f"{d['BAIXADOS']} posts do Mastodon · agosto de 2026",
          tamanho=12.5, cor=MUDO, espacamento=1.4)
    return slide


def slide_algoritmo(apresentacao):
    slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])
    fundo_papel(slide)
    cabecalho(slide,
              "Como decidimos se um post fala bem ou mal da Argentina",
              "Cada palavra do dicionário ANEW já vem com uma nota de agrado, dada por "
              "voluntários. A nota de um post é a média das notas das suas palavras.")

    col = Inches(5.9)
    y = Inches(2.15)
    passo(slide, MARGEM, y, col, 1, "Buscamos os posts",
          "Um programa nosso baixa da API do Mastodon os posts que citam "
          "**Argentina**, e confere o idioma de cada um.")
    passo(slide, MARGEM, y + Inches(1.15), col, 2, "Casamos as palavras com o dicionário",
          "As 1.034 palavras do ANEW têm nota de 0 a 100. **love** = 99 · "
          "**news** = 60 · **terrible** = 22. Reduzimos cada palavra do post à "
          "forma do dicionário, para que _loved_ encontre _love_.")
    passo(slide, MARGEM, y + Inches(2.55), col, 3, "Tiramos a média e comparamos",
          "Acima da média do dicionário é positivo, abaixo é negativo, perto "
          "dela é neutro.")

    dir_x = MARGEM + Inches(6.5)
    dir_w = Inches(5.45)

    cartao(slide, dir_x, y, dir_w, Inches(1.5))
    texto(slide, dir_x + Inches(0.3), y + Inches(0.3), dir_w - Inches(0.6), Inches(0.4),
          "nota do post = média das notas das suas palavras",
          tamanho=15, cor=TINTA, alinhamento=PP_ALIGN.CENTER)
    texto(slide, dir_x + Inches(0.3), y + Inches(0.82), dir_w - Inches(0.6), Inches(0.6),
          "“Argentina is **beautiful**, I **love** the **people** and the **food**”\n"
          "(86 + 99 + 83 + 87) ÷ 4 = **89** → positivo",
          tamanho=11, cor=MUDO, alinhamento=PP_ALIGN.CENTER, espacamento=1.35)

    cartao(slide, dir_x, y + Inches(1.85), dir_w, Inches(2.35), acento=True)
    texto(slide, dir_x + Inches(0.32), y + Inches(2.1), dir_w - Inches(0.6), Inches(0.25),
          "AS TRÊS DIMENSÕES DO ANEW", tamanho=9, cor=MUDO, negrito=True)
    texto(slide, dir_x + Inches(0.32), y + Inches(2.45), dir_w - Inches(0.62), Inches(1.3),
          "O dicionário mede **Pleasure** (agradável ou desagradável), "
          "**Arousal** (calmo ou agitado) e **Dominance** (quanto controle a "
          "palavra transmite).\n"
          "Usamos só **Pleasure**, porque é a dimensão que responde “falaram bem "
          "ou mal?”. As outras duas medem intensidade e poder, não aprovação.",
          tamanho=11.5, cor=TINTA2, espacamento=1.3)

    rodape(slide, "Dicionário ANEW (Bradley & Lang, 1999) · entidade: Argentina · "
                  "fonte: Mastodon", 1)
    return slide


def slide_escolhas(apresentacao, d):
    slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])
    fundo_papel(slide)
    cabecalho(slide, "Por que escolhemos assim",
              "O enunciado deixa em aberto como transformar notas de palavras numa "
              "decisão sobre o post. Estas foram as nossas escolhas, e o motivo de cada uma.",
              fraco="  — três decisões que mudam o resultado")

    y = Inches(2.05)
    esq_x, esq_w = MARGEM, Inches(5.85)

    # escolha 1
    cartao(slide, esq_x, y, esq_w, Inches(1.6), acento=True)
    texto(slide, esq_x + Inches(0.3), y + Inches(0.2), esq_w - Inches(0.55), Inches(0.2),
          "ESCOLHA 1 · UMA ENTIDADE SÓ", tamanho=9, cor=MUDO, negrito=True)
    texto(slide, esq_x + Inches(0.3), y + Inches(0.45), esq_w - Inches(0.55), Inches(0.25),
          "Buscamos apenas “Argentina”", tamanho=14, cor=TINTA, negrito=True)
    texto(slide, esq_x + Inches(0.3), y + Inches(0.75), esq_w - Inches(0.58), Inches(0.5),
          "Testamos buscar também Buenos Aires e Patagônia, para ter mais posts. "
          "Mas isso trazia turismo e natureza — **outras entidades** — e empurrava "
          "o resultado para o positivo.",
          tamanho=11, cor=TINTA2, espacamento=1.28)

    # escolha 2
    y2 = y + Inches(1.78)
    cartao(slide, esq_x, y2, esq_w, Inches(2.6), acento=True)
    texto(slide, esq_x + Inches(0.3), y2 + Inches(0.2), esq_w - Inches(0.55), Inches(0.2),
          "ESCOLHA 2 · O PONTO DE CORTE", tamanho=9, cor=MUDO, negrito=True)
    texto(slide, esq_x + Inches(0.3), y2 + Inches(0.45), esq_w - Inches(0.55), Inches(0.25),
          "Comparamos com 58, não com 50", tamanho=14, cor=TINTA, negrito=True)
    texto(slide, esq_x + Inches(0.3), y2 + Inches(0.75), esq_w - Inches(0.58), Inches(0.5),
          "A escala vai de 0 a 100, então o meio _parece_ ser 50. Mas o ANEW não é "
          "equilibrado: **62% das suas palavras estão acima de 50**, e a média delas é 58.",
          tamanho=11, cor=TINTA2, espacamento=1.28)

    barra_w = esq_w - Inches(0.6)
    texto(slide, esq_x + Inches(0.3), y2 + Inches(1.42), barra_w, Inches(0.18),
          "os mesmos posts, cortando em 50", tamanho=9, cor=MUDO)
    barra_empilhada(slide, esq_x + Inches(0.3), y2 + Inches(1.62), barra_w, Inches(0.26),
                    [(d["ING_NEG"], NEG, None),
                     (d["ING_POS"], POS, f"{d['ING_POS']}% positivo")])
    texto(slide, esq_x + Inches(0.3), y2 + Inches(2.0), barra_w, Inches(0.18),
          "cortando em 58, a média real do dicionário", tamanho=9, cor=MUDO)
    barra_empilhada(slide, esq_x + Inches(0.3), y2 + Inches(2.2), barra_w, Inches(0.26),
                    [(d["PCT_NEG"], NEG, None),
                     (d["PCT_NEU"], NEU, f"{d['PCT_NEU']}% neutro"),
                     (d["PCT_POS"], POS, f"{d['PCT_POS']}% positivo")])

    # escolha 3 + o caso
    dir_x, dir_w = MARGEM + Inches(6.5), Inches(5.45)
    cartao(slide, dir_x, y, dir_w, Inches(1.05), acento=True)
    texto(slide, dir_x + Inches(0.3), y + Inches(0.18), dir_w - Inches(0.55), Inches(0.2),
          "ESCOLHA 3 · QUANDO NÃO OPINAR", tamanho=9, cor=MUDO, negrito=True)
    texto(slide, dir_x + Inches(0.3), y + Inches(0.42), dir_w - Inches(0.58), Inches(0.5),
          "Post com menos de 3 palavras conhecidas é descartado",
          tamanho=14, cor=TINTA, negrito=True, espacamento=1.15)

    yc = y + Inches(1.25)
    cartao(slide, dir_x, yc, dir_w, Inches(2.95))
    texto(slide, dir_x + Inches(0.28), yc + Inches(0.2), dir_w - Inches(0.55), Inches(0.2),
          "POST REAL · O ANEW SÓ RECONHECEU UMA PALAVRA", tamanho=9, cor=MUDO, negrito=True)
    texto(slide, dir_x + Inches(0.28), yc + Inches(0.48), dir_w - Inches(0.56), Inches(0.95),
          "_“An elderly lady of the indigenous Humahuaca community in Jujuy cries "
          "as the puppet regime in Argentina evicts her from her _**home · 90**_ "
          "and land. Large-scale protests against the law continue.”_",
          tamanho=11, cor=TINTA, espacamento=1.32)

    meio = dir_w / 2 - Inches(0.32)
    yd = yc + Inches(1.62)
    forma(slide, MSO_SHAPE.ROUNDED_RECTANGLE, dir_x + Inches(0.28), yd, meio, Inches(0.72),
          fundo=NEG_FUNDO, borda=NEG)
    texto(slide, dir_x + Inches(0.42), yd + Inches(0.1), meio - Inches(0.28), Inches(0.2),
          "sem regra nenhuma:", tamanho=9.5, cor=TINTA2)
    texto(slide, dir_x + Inches(0.42), yd + Inches(0.33), meio - Inches(0.28), Inches(0.28),
          "POSITIVO · 90", tamanho=14, cor=NEG, negrito=True)

    dx2 = dir_x + Inches(0.28) + meio + Inches(0.18)
    forma(slide, MSO_SHAPE.ROUNDED_RECTANGLE, dx2, yd, meio, Inches(0.72),
          fundo=POS_FUNDO, borda=POS)
    texto(slide, dx2 + Inches(0.14), yd + Inches(0.1), meio - Inches(0.28), Inches(0.2),
          "com a nossa regra:", tamanho=9.5, cor=TINTA2)
    texto(slide, dx2 + Inches(0.14), yd + Inches(0.33), meio - Inches(0.28), Inches(0.28),
          "SEM SINAL", tamanho=14, cor=POS, negrito=True)

    texto(slide, dir_x + Inches(0.28), yd + Inches(0.9), dir_w - Inches(0.56), Inches(0.5),
          "De 63 palavras, o dicionário só reconheceu **home**. Os termos que dariam "
          "o tom — _cries_, _evicts_, _protests_ — não estão no ANEW.",
          tamanho=10.5, cor=TINTA2, espacamento=1.28)

    rodape(slide, "Cada escolha foi testada contra os posts reais antes de ser adotada", 2)
    return slide


def slide_resultado(apresentacao, d):
    slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])
    fundo_papel(slide)
    cabecalho(slide, "Fala-se mais bem que mal da Argentina",
              f"{d['CLASSIFICADOS']} posts classificados, de {d['INGLES']} em inglês. "
              f"Mais da metade é positiva e quatro em cada dez ficam no meio da escala. "
              f"Crítica forte é rara: {d['PCT_NEG']}%.")

    figura = RAIZ / "figuras" / "sentimento.png"
    largura = Inches(10.2)
    x = (Inches(13.333) - largura) // 2
    slide.shapes.add_picture(str(figura), x, Inches(1.9), width=largura)

    rodape(slide, f"{d['BAIXADOS']} posts baixados · {d['INGLES']} confirmados em inglês · "
                  f"{d['CLASSIFICADOS']} com palavras suficientes para pontuar", 3)
    return slide


def main():
    d = json.load(open(AQUI / "numeros.json"))

    apresentacao = Presentation()
    apresentacao.slide_width = Inches(13.333)
    apresentacao.slide_height = Inches(7.5)

    capa(apresentacao, d)
    slide_algoritmo(apresentacao)
    slide_escolhas(apresentacao, d)
    slide_resultado(apresentacao, d)

    apresentacao.save(SAIDA)
    print(f"  {len(apresentacao.slides)} slides gravados em "
          f"{SAIDA.relative_to(RAIZ)}  (texto editável)")


if __name__ == "__main__":
    main()
